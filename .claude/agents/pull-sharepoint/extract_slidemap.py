"""
Pass 1: PPTX Slide Map Extraction
Used by pull-sharepoint agent to classify slides before visual analysis.

Usage:
    python extract_slidemap.py <pptx_path> [--output <json_path>]

Output:
    - Prints slide map summary to stdout
    - Writes full slide map JSON to --output path (or stdout if not specified)
    - Writes full text content per slide to <output>.texts.json
"""

import sys
import json
import argparse
from pathlib import Path
from pptx import Presentation


VISUAL_KEYWORDS = [
    'Architecture', 'Diagram', 'Overview', 'Workflow', 'Matrix', 'Flow',
    'Org Chart', 'Gantt', 'Roadmap', 'Components', 'Integration', 'Landscape'
]


def classify_slide(word_count, shapes_without_text, has_table, title):
    """Classify a slide based on text/shape metrics."""
    wc = word_count
    swot = shapes_without_text

    if wc < 5 and swot == 0:
        cls = 'blank'
    elif wc < 15 and swot >= 1:
        cls = 'visual-dominant'
    elif wc >= 30 and swot <= 2:
        cls = 'text-rich'
    elif wc >= 15 and swot >= 1:
        cls = 'mixed'
    elif has_table:
        cls = 'table-only'
    else:
        cls = 'text-rich'

    # Title keyword upgrade
    if any(kw.lower() in title.lower() for kw in VISUAL_KEYWORDS):
        if cls in ('text-rich', 'mixed'):
            cls = 'visual-dominant'

    return cls


def detect_template_tail(slide_map):
    """Detect template/master slides at the end of a deck.

    Pattern: A run of consecutive slides at the tail with very low word count
    and repetitive minimal content (page numbers, 'CONFIDENTIAL', empty placeholders).
    Returns the first slide number where the template tail begins, or None.
    """
    if len(slide_map) < 10:
        return None

    # Walk backwards from end, looking for consecutive low-content slides
    MIN_TAIL_RUN = 5  # need at least 5 consecutive template slides to trigger
    TEMPLATE_WORD_THRESHOLD = 10  # slides with <= this many words are candidate templates

    tail_start = None
    consecutive = 0

    for i in range(len(slide_map) - 1, -1, -1):
        s = slide_map[i]
        if s['word_count'] <= TEMPLATE_WORD_THRESHOLD:
            consecutive += 1
            tail_start = i
        else:
            break

    if consecutive >= MIN_TAIL_RUN:
        return slide_map[tail_start]['slide_num']
    return None


def extract_slidemap(pptx_path):
    """Extract slide map with classification and full text from a PPTX file."""
    prs = Presentation(pptx_path)
    slide_map = []

    for slide_num, slide in enumerate(prs.slides, 1):
        title = ''
        word_count = 0
        shapes_with_text = 0
        shapes_without_text = 0
        has_table = False
        texts = []

        for shape in slide.shapes:
            if shape.has_table:
                has_table = True
                table_texts = []
                for row in shape.table.rows:
                    row_cells = []
                    for cell in row.cells:
                        t = cell.text.strip()
                        if t:
                            row_cells.append(t)
                            word_count += len(t.split())
                    if row_cells:
                        table_texts.append(' | '.join(row_cells))
                if table_texts:
                    texts.append('[TABLE]\n' + '\n'.join(table_texts))
                shapes_with_text += 1
            elif shape.has_text_frame:
                shape_text = ''
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        shape_text += t + '\n'
                if shape_text.strip():
                    texts.append(shape_text.strip())
                    word_count += len(shape_text.split())
                    shapes_with_text += 1
                    # Check if title placeholder
                    if shape.is_placeholder and shape.placeholder_format.idx == 0:
                        title = shape_text.strip().split('\n')[0]
                else:
                    shapes_without_text += 1
            else:
                shapes_without_text += 1

        # Fallback title
        if not title and texts:
            title = texts[0].split('\n')[0][:80]

        classification = classify_slide(word_count, shapes_without_text, has_table, title)

        slide_map.append({
            'slide_num': slide_num,
            'title': title,
            'word_count': word_count,
            'shapes_with_text': shapes_with_text,
            'shapes_without_text': shapes_without_text,
            'has_table': has_table,
            'classification': classification,
            'texts': texts
        })

    return slide_map


def print_summary(slide_map, template_tail_start=None):
    """Print slide map summary to stderr (so stdout stays clean for JSON)."""
    counts = {}
    for s in slide_map:
        counts[s['classification']] = counts.get(s['classification'], 0) + 1

    print(f"SLIDE_MAP:", file=sys.stderr)
    print(f"Total slides: {len(slide_map)}", file=sys.stderr)
    for cls in ['text-rich', 'mixed', 'visual-dominant', 'table-only', 'blank', 'template']:
        if counts.get(cls, 0) > 0:
            print(f"  {cls:16s}: {counts.get(cls, 0):3d} slides", file=sys.stderr)

    visual_slides = [s['slide_num'] for s in slide_map
                     if s['classification'] in ('visual-dominant', 'mixed')]
    print(f"Visual pass needed: {len(visual_slides)} slides", file=sys.stderr)
    print(f"  Slides: {visual_slides}", file=sys.stderr)

    if template_tail_start:
        print(f"Template tail detected: slides {template_tail_start}+ (skipped)", file=sys.stderr)


def main():
    parser = argparse.ArgumentParser(description='Extract slide map from PPTX')
    parser.add_argument('pptx_path', help='Path to the PPTX file')
    parser.add_argument('--output', '-o', help='Output JSON path (slide map without texts)')
    args = parser.parse_args()

    pptx_path = Path(args.pptx_path)
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    slide_map = extract_slidemap(pptx_path)

    # Detect and reclassify template tail
    template_tail_start = detect_template_tail(slide_map)
    if template_tail_start:
        for s in slide_map:
            if s['slide_num'] >= template_tail_start:
                s['classification'] = 'template'

    print_summary(slide_map, template_tail_start)

    # Separate texts from map metadata
    map_meta = [{k: v for k, v in s.items() if k != 'texts'} for s in slide_map]
    texts_data = [{'slide_num': s['slide_num'], 'title': s['title'],
                   'classification': s['classification'], 'texts': s['texts']}
                  for s in slide_map]

    if args.output:
        output_path = Path(args.output)
        # Write map (without texts)
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(map_meta, f, indent=2, ensure_ascii=False)
        # Write texts separately
        texts_path = output_path.with_suffix('.texts.json')
        with open(texts_path, 'w', encoding='utf-8') as f:
            json.dump(texts_data, f, indent=2, ensure_ascii=False)
        print(f"\nMap written to: {output_path}", file=sys.stderr)
        print(f"Texts written to: {texts_path}", file=sys.stderr)
    else:
        # Output map JSON to stdout
        json.dump(map_meta, sys.stdout, indent=2, ensure_ascii=False)


if __name__ == '__main__':
    main()
